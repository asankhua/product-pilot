"""
PPT Generator Core Module
Handles professional PowerPoint creation with templates, charts, and layouts
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import ChartData

import matplotlib.pyplot as plt
import matplotlib
matplotlib.use('Agg')  # Non-interactive backend
from io import BytesIO
import os
import tempfile
from typing import List, Optional, Dict, Any
import logging

logger = logging.getLogger(__name__)


class PPTGenerator:
    """Professional PowerPoint presentation generator"""
    
    # Color schemes for different templates
    COLOR_SCHEMES = {
        "professional": {
            "primary": RGBColor(30, 64, 175),      # Blue 800
            "secondary": RGBColor(59, 130, 246),   # Blue 500
            "accent": RGBColor(16, 185, 129),      # Emerald 500
            "text": RGBColor(31, 41, 55),        # Gray 800
            "light": RGBColor(243, 244, 246),    # Gray 100
            "white": RGBColor(255, 255, 255),
            "dark": RGBColor(17, 24, 39),        # Gray 900
        },
        "minimal": {
            "primary": RGBColor(31, 41, 55),       # Gray 800
            "secondary": RGBColor(107, 114, 128), # Gray 500
            "accent": RGBColor(99, 102, 241),     # Indigo 500
            "text": RGBColor(31, 41, 55),
            "light": RGBColor(249, 250, 251),    # Gray 50
            "white": RGBColor(255, 255, 255),
            "dark": RGBColor(17, 24, 39),
        },
        "dark": {
            "primary": RGBColor(99, 102, 241),     # Indigo 500
            "secondary": RGBColor(139, 92, 246),   # Violet 500
            "accent": RGBColor(236, 72, 153),      # Pink 500
            "text": RGBColor(255, 255, 255),
            "light": RGBColor(55, 65, 81),        # Gray 700
            "white": RGBColor(255, 255, 255),
            "dark": RGBColor(17, 24, 39),         # Gray 900
        },
        "startup": {
            "primary": RGBColor(244, 63, 94),      # Rose 500
            "secondary": RGBColor(251, 146, 60),   # Orange 400
            "accent": RGBColor(234, 179, 8),      # Yellow 500
            "text": RGBColor(31, 41, 55),
            "light": RGBColor(255, 241, 242),     # Rose 50
            "white": RGBColor(255, 255, 255),
            "dark": RGBColor(136, 19, 55),        # Rose 900
        }
    }
    
    def __init__(self):
        self.template_dir = os.path.join(os.path.dirname(__file__), "..", "templates")
    
    def create_presentation(
        self,
        project_name: str,
        project_description: Optional[str],
        steps: List[Any],
        template: str,
        output_path: str,
        include_charts: bool = True,
        include_images: bool = True,
        preview_mode: bool = False
    ):
        """Create a complete presentation"""
        
        colors = self.COLOR_SCHEMES.get(template, self.COLOR_SCHEMES["professional"])
        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
        prs.slide_height = Inches(7.5)
        
        # Set presentation properties
        prs.core_properties.title = f"{project_name} - Product Strategy"
        prs.core_properties.author = "Product Pilot"
        prs.core_properties.subject = "Product Strategy Presentation"
        
        # Title Slide
        self._add_title_slide(prs, project_name, project_description, colors)
        
        if preview_mode:
            prs.save(output_path)
            return
        
        # Agenda Slide
        self._add_agenda_slide(prs, steps, colors)
        
        # Content Slides
        for step in steps:
            self._add_step_slide(prs, step, colors, include_charts)
        
        # Thank You Slide
        self._add_closing_slide(prs, colors)
        
        # Save presentation
        prs.save(output_path)
        logger.info(f"Saved presentation to: {output_path}")
    
    def _add_title_slide(self, prs: Presentation, title: str, subtitle: Optional[str], colors: Dict):
        """Create professional title slide"""
        blank_slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Background shape (full slide gradient effect)
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), 
            prs.slide_width, prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = colors["light"]
        background.line.fill.background()
        
        # Top accent bar
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
            prs.slide_width, Inches(0.3)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = colors["primary"]
        accent_bar.line.fill.background()
        
        # Decorative circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(11), Inches(5.5), Inches(2), Inches(2)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = colors["secondary"]
        circle.fill.fore_color.brightness = 0.3
        circle.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(10), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(54)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = colors["primary"]
        title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        
        # Subtitle
        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(0.8))
            sub_frame = sub_box.text_frame
            sub_frame.text = subtitle
            sub_frame.paragraphs[0].font.size = Pt(24)
            sub_frame.paragraphs[0].font.color.rgb = colors["secondary"]
            sub_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        
        # Date and author
        meta_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(6), Inches(0.5))
        meta_frame = meta_box.text_frame
        meta_frame.text = f"Generated by Product Pilot • {self._get_date()}"
        meta_frame.paragraphs[0].font.size = Pt(12)
        meta_frame.paragraphs[0].font.color.rgb = colors["text"]
        meta_frame.paragraphs[0].font.italic = True
    
    def _add_agenda_slide(self, prs: Presentation, steps: List[Any], colors: Dict):
        """Create agenda/overview slide"""
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Header bar
        self._add_header_bar(slide, "Presentation Overview", colors, prs.slide_width)
        
        # Agenda items
        y_pos = 1.8
        for i, step in enumerate(steps[:6]):  # Max 6 items
            step_num = i + 1
            step_name = step.stepName if hasattr(step, 'stepName') else f"Step {step_num}"
            
            # Number circle
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(1), Inches(y_pos), Inches(0.5), Inches(0.5)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = colors["primary"]
            circle.line.fill.background()
            
            # Number text
            num_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos + 0.1), Inches(0.5), Inches(0.3))
            num_frame = num_box.text_frame
            num_frame.text = str(step_num)
            num_frame.paragraphs[0].font.size = Pt(14)
            num_frame.paragraphs[0].font.bold = True
            num_frame.paragraphs[0].font.color.rgb = colors["white"]
            num_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Step name
            name_box = slide.shapes.add_textbox(Inches(1.8), Inches(y_pos + 0.05), Inches(6), Inches(0.4))
            name_frame = name_box.text_frame
            name_frame.text = step_name
            name_frame.paragraphs[0].font.size = Pt(18)
            name_frame.paragraphs[0].font.color.rgb = colors["text"]
            
            y_pos += 0.8
    
    def _add_step_slide(self, prs: Presentation, step: Any, colors: Dict, include_charts: bool):
        """Add a content slide for a pipeline step"""
        step_id = step.stepId if hasattr(step, 'stepId') else 0
        step_name = step.stepName if hasattr(step, 'stepName') else "Step"
        data = step.data if hasattr(step, 'data') else {}
        
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Header bar with step name
        self._add_header_bar(slide, f"{step_id}. {step_name}", colors, prs.slide_width)
        
        # Route to specific slide type based on step
        if step_id == 1:
            self._add_reframe_slide(slide, data, colors)
        elif step_id == 2:
            self._add_vision_slide(slide, data, colors)
        elif step_id == 3:
            self._add_personas_slide(slide, data, colors)
        elif step_id == 4:
            self._add_questions_slide(slide, data, colors)
        elif step_id == 5:
            self._add_market_slide(slide, data, colors, include_charts)
        elif step_id == 6:
            self._add_prd_slide(slide, data, colors)
        elif step_id == 7:
            self._add_stories_slide(slide, data, colors)
        elif step_id == 8:
            self._add_roadmap_slide(slide, data, colors)
        elif step_id == 9:
            self._add_okrs_slide(slide, data, colors, include_charts)
        else:
            self._add_generic_slide(slide, data, colors)
    
    def _add_header_bar(self, slide, title: str, colors: Dict, slide_width):
        """Add consistent header bar to slides"""
        # Header background
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
            slide_width, Inches(1)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = colors["primary"]
        header.line.fill.background()
        
        # Title text
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = colors["white"]
    
    def _add_reframe_slide(self, slide, data: Dict, colors: Dict):
        """Reframe problem statement slide"""
        y_pos = 1.3
        
        # Problem title
        if data.get('problemTitle'):
            self._add_text_box(slide, "Problem Title:", data['problemTitle'], Inches(0.5), Inches(y_pos), colors, bold_label=True)
            y_pos += 1.0
        
        # Reframed problem
        if data.get('reframedProblem'):
            self._add_text_box(slide, "Reframed Problem:", data['reframedProblem'], Inches(0.5), Inches(y_pos), colors, bold_label=True)
            y_pos += 1.2
        
        # Root causes
        if data.get('rootCauses'):
            self._add_bullet_list(slide, "Root Causes:", data['rootCauses'], Inches(0.5), Inches(y_pos), Inches(6), colors)
        
        # Right column - Impact & Opportunity
        if data.get('userImpact') or data.get('opportunitySize'):
            x_right = Inches(7)
            y_right = 1.3
            
            if data.get('userImpact'):
                self._add_info_card(slide, "User Impact", data['userImpact'], x_right, Inches(y_right), Inches(5.5), Inches(1.5), colors["accent"])
                y_right += 1.8
            
            if data.get('opportunitySize'):
                self._add_info_card(slide, "Market Opportunity", data['opportunitySize'], x_right, Inches(y_right), Inches(5.5), Inches(1.5), colors["secondary"])
    
    def _add_vision_slide(self, slide, data: Dict, colors: Dict):
        """Product vision slide"""
        # Vision statement (prominent)
        if data.get('visionStatement'):
            vision_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.3), Inches(12), Inches(1.2)
            )
            vision_box.fill.solid()
            vision_box.fill.fore_color.rgb = colors["light"]
            vision_box.line.color.rgb = colors["primary"]
            
            text_frame = vision_box.text_frame
            text_frame.text = data['visionStatement']
            text_frame.paragraphs[0].font.size = Pt(20)
            text_frame.paragraphs[0].font.italic = True
            text_frame.paragraphs[0].font.color.rgb = colors["primary"]
            text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            text_frame.word_wrap = True
        
        # Two column layout for other details
        y_pos = 2.8
        
        # Left column
        if data.get('elevatorPitch'):
            self._add_text_box(slide, "Elevator Pitch:", data['elevatorPitch'], Inches(0.5), Inches(y_pos), colors)
            y_pos += 1.5
        
        if data.get('targetAudience'):
            self._add_text_box(slide, "Target Audience:", data['targetAudience'], Inches(0.5), Inches(y_pos), colors)
        
        # Right column
        x_right = Inches(6.5)
        y_right = 2.8
        
        if data.get('valueProposition'):
            self._add_info_card(slide, "Value Proposition", data['valueProposition'], Inches(x_right), Inches(y_right), Inches(6), Inches(1.5), colors["accent"])
            y_right += 1.8
        
        # Success metrics
        if data.get('successMetrics'):
            self._add_bullet_list(slide, "Success Metrics:", data['successMetrics'][:4], x_right, Inches(y_right), Inches(6), colors)
    
    def _add_personas_slide(self, slide, data: Dict, colors: Dict):
        """User personas slide with cards"""
        personas = data.get('personas', [])
        if not personas and isinstance(data, list):
            personas = data
        
        if not personas:
            self._add_text_box(slide, "", "No persona data available", Inches(0.5), Inches(1.5), colors)
            return
        
        # Layout: up to 3 persona cards side by side
        card_width = Inches(4)
        gap = Inches(0.3)
        start_x = Inches(0.5)
        y_pos = Inches(1.3)
        
        for i, persona in enumerate(personas[:3]):
            x_pos = start_x + (i * (card_width + gap))
            self._add_persona_card(slide, persona, x_pos, y_pos, card_width, colors)
    
    def _add_persona_card(self, slide, persona: Dict, x, y, width, colors):
        """Create a persona info card"""
        # Card background
        card_height = Inches(5.5)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, card_height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = colors["light"]
        card.line.color.rgb = colors["secondary"]
        
        # Persona name & role
        name_y = y + Inches(0.2)
        name_box = slide.shapes.add_textbox(x + Inches(0.2), name_y, width - Inches(0.4), Inches(0.5))
        name_frame = name_box.text_frame
        name_frame.text = f"{persona.get('name', 'Unknown')} - {persona.get('role', 'User')}"
        name_frame.paragraphs[0].font.size = Pt(16)
        name_frame.paragraphs[0].font.bold = True
        name_frame.paragraphs[0].font.color.rgb = colors["primary"]
        
        # Bio
        bio_y = name_y + Inches(0.6)
        bio_text = persona.get('bio', '')[:150] + "..." if len(persona.get('bio', '')) > 150 else persona.get('bio', '')
        self._add_wrapped_text(slide, bio_text, x + Inches(0.2), bio_y, width - Inches(0.4), Pt(10), colors["text"])
        
        # Pain points
        if persona.get('painPoints'):
            pain_y = bio_y + Inches(1.0)
            pain_box = slide.shapes.add_textbox(x + Inches(0.2), pain_y, width - Inches(0.4), Inches(0.3))
            pain_frame = pain_box.text_frame
            pain_frame.text = "Pain Points:"
            pain_frame.paragraphs[0].font.size = Pt(11)
            pain_frame.paragraphs[0].font.bold = True
            pain_frame.paragraphs[0].font.color.rgb = RGBColor(239, 68, 68)  # Red
            
            for j, pain in enumerate(persona['painPoints'][:2]):
                pain_item_y = pain_y + Inches(0.25 + (j * 0.25))
                pain_item = slide.shapes.add_textbox(x + Inches(0.3), pain_item_y, width - Inches(0.5), Inches(0.25))
                pain_item_frame = pain_item.text_frame
                pain_item_frame.text = f"• {pain[:50]}"
                pain_item_frame.paragraphs[0].font.size = Pt(9)
                pain_item_frame.paragraphs[0].font.color.rgb = colors["text"]
        
        # Goals
        if persona.get('goals'):
            goal_y = y + Inches(3.8)
            goal_box = slide.shapes.add_textbox(x + Inches(0.2), goal_y, width - Inches(0.4), Inches(0.3))
            goal_frame = goal_box.text_frame
            goal_frame.text = "Goals:"
            goal_frame.paragraphs[0].font.size = Pt(11)
            goal_frame.paragraphs[0].font.bold = True
            goal_frame.paragraphs[0].font.color.rgb = colors["accent"]
            
            for j, goal in enumerate(persona['goals'][:2]):
                goal_item_y = goal_y + Inches(0.25 + (j * 0.25))
                goal_item = slide.shapes.add_textbox(x + Inches(0.3), goal_item_y, width - Inches(0.5), Inches(0.25))
                goal_item_frame = goal_item.text_frame
                goal_item_frame.text = f"• {goal[:50]}"
                goal_item_frame.paragraphs[0].font.size = Pt(9)
                goal_item_frame.paragraphs[0].font.color.rgb = colors["text"]
    
    def _add_questions_slide(self, slide, data: Dict, colors: Dict):
        """Q&A slide"""
        questions = data.get('questions', [])
        if not questions and isinstance(data, list):
            questions = data
        
        if not questions:
            self._add_text_box(slide, "", "No questions data available", Inches(0.5), Inches(1.5), colors)
            return
        
        # Two column layout
        left_x = Inches(0.5)
        right_x = Inches(6.5)
        y_start = 1.3
        y_left = y_start
        y_right = y_start
        
        for i, q in enumerate(questions[:6]):  # Max 6 questions
            x = left_x if i % 2 == 0 else right_x
            y = y_left if i % 2 == 0 else y_right
            
            # Question card
            q_text = q.get('question', '')[:80]
            a_text = q.get('aiAnswer', q.get('userAnswer', ''))[:100]
            
            # Question
            q_box = slide.shapes.add_textbox(x, Inches(y), Inches(6), Inches(0.5))
            q_frame = q_box.text_frame
            q_frame.text = f"Q{i+1}: {q_text}"
            q_frame.paragraphs[0].font.size = Pt(12)
            q_frame.paragraphs[0].font.bold = True
            q_frame.paragraphs[0].font.color.rgb = colors["primary"]
            
            # Answer
            if a_text:
                a_box = slide.shapes.add_textbox(x + Inches(0.2), Inches(y + 0.4), Inches(5.8), Inches(0.6))
                a_frame = a_box.text_frame
                a_frame.text = a_text
                a_frame.paragraphs[0].font.size = Pt(10)
                a_frame.paragraphs[0].font.color.rgb = colors["text"]
                a_frame.word_wrap = True
            
            if i % 2 == 0:
                y_left += 1.2
            else:
                y_right += 1.2
    
    def _add_market_slide(self, slide, data: Dict, colors: Dict, include_charts: bool):
        """Market analysis slide with TAM/SAM/SOM chart"""
        # Market overview
        if data.get('marketOverview'):
            self._add_text_box(slide, "Market Overview:", data['marketOverview'][:200], Inches(0.5), Inches(1.3), colors)
        
        # Competitors table
        competitors = data.get('competitors', [])
        if competitors:
            self._add_competitor_table(slide, competitors[:4], Inches(0.5), Inches(3.0), Inches(7), colors)
        
        # TAM/SAM/SOM visualization
        if data.get('tam') or data.get('sam') or data.get('som'):
            if include_charts:
                self._add_tam_sam_som_chart(slide, data, Inches(8), Inches(1.5), colors)
            else:
                # Text version
                x = Inches(8)
                y = 1.5
                for label, key in [("TAM", "tam"), ("SAM", "sam"), ("SOM", "som")]:
                    if data.get(key):
                        self._add_info_card(slide, label, data[key], Inches(x), Inches(y), Inches(4.5), Inches(0.8), colors["secondary"])
                        y += 1.0
    
    def _add_tam_sam_som_chart(self, slide, data: Dict, x, y, colors):
        """Add TAM/SAM/SOM funnel chart using matplotlib"""
        try:
            # Create funnel chart
            fig, ax = plt.subplots(figsize=(4, 3))
            
            labels = []
            values = []
            for label, key in [("TAM", "tam"), ("SAM", "sam"), ("SOM", "som")]:
                if data.get(key):
                    labels.append(label)
                    # Try to extract numeric value or use placeholder
                    val_str = str(data[key])
                    values.append(len(val_str))  # Size based on content length for visual effect
            
            if values:
                colors_hex = ["#1E40AF", "#3B82F6", "#60A5FA"]
                ax.barh(labels[::-1], values[::-1], color=colors_hex[:len(values)])
                ax.set_xlabel("Market Size")
                ax.set_title("Market Opportunity")
                plt.tight_layout()
                
                # Save to buffer
                buffer = BytesIO()
                plt.savefig(buffer, format='png', dpi=100, bbox_inches='tight')
                buffer.seek(0)
                
                # Add to slide
                slide.shapes.add_picture(buffer, x, y, width=Inches(4.5))
                plt.close(fig)
        except Exception as e:
            logger.warning(f"Failed to create chart: {e}")
    
    def _add_competitor_table(self, slide, competitors, x, y, width, colors):
        """Add competitor analysis table"""
        rows = len(competitors) + 1
        cols = 3
        
        table = slide.shapes.add_table(rows, cols, x, y, width, Inches(0.6 * rows)).table
        
        # Header
        headers = ["Competitor", "Strengths", "Weaknesses"]
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = colors["primary"]
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(10)
            paragraph.font.bold = True
            paragraph.font.color.rgb = colors["white"]
        
        # Data rows
        for i, comp in enumerate(competitors):
            row = i + 1
            
            # Name
            cell = table.cell(row, 0)
            cell.text = comp.get('name', 'Unknown')
            cell.text_frame.paragraphs[0].font.size = Pt(10)
            
            # Strengths
            cell = table.cell(row, 1)
            strengths = comp.get('strengths', [])
            cell.text = "\n".join([f"• {s[:30]}" for s in strengths[:2]])
            cell.text_frame.paragraphs[0].font.size = Pt(9)
            
            # Weaknesses
            cell = table.cell(row, 2)
            weaknesses = comp.get('weaknesses', [])
            cell.text = "\n".join([f"• {w[:30]}" for w in weaknesses[:2]])
            cell.text_frame.paragraphs[0].font.size = Pt(9)
    
    def _add_prd_slide(self, slide, data: Dict, colors: Dict):
        """PRD slide with feature list"""
        features = data.get('features', [])
        if not features:
            # Try to extract from saved data
            if isinstance(data, dict):
                for key in ['featureList', 'capabilities', 'requirements']:
                    if key in data and isinstance(data[key], list):
                        features = data[key]
                        break
        
        if not features:
            self._add_text_box(slide, "PRD Content", str(data)[:300], Inches(0.5), Inches(1.5), colors)
            return
        
        # Feature list with priority indicators
        y_pos = 1.5
        for i, feature in enumerate(features[:8]):  # Max 8 features
            feature_text = feature.get('name', feature.get('title', str(feature)))[:60]
            priority = feature.get('priority', 'Medium')
            
            # Priority color
            if priority.lower() == 'high':
                p_color = RGBColor(239, 68, 68)  # Red
            elif priority.lower() == 'medium':
                p_color = RGBColor(245, 158, 11)  # Orange
            else:
                p_color = RGBColor(16, 185, 129)  # Green
            
            # Priority badge
            badge = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y_pos), Inches(0.8), Inches(0.3)
            )
            badge.fill.solid()
            badge.fill.fore_color.rgb = p_color
            badge.line.fill.background()
            
            badge_text = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos + 0.05), Inches(0.8), Inches(0.2))
            badge_text_frame = badge_text.text_frame
            badge_text_frame.text = priority[:3].upper()
            badge_text_frame.paragraphs[0].font.size = Pt(9)
            badge_text_frame.paragraphs[0].font.bold = True
            badge_text_frame.paragraphs[0].font.color.rgb = colors["white"]
            badge_text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Feature text
            feature_box = slide.shapes.add_textbox(Inches(1.5), Inches(y_pos), Inches(11), Inches(0.3))
            feature_frame = feature_box.text_frame
            feature_frame.text = feature_text
            feature_frame.paragraphs[0].font.size = Pt(12)
            feature_frame.paragraphs[0].font.color.rgb = colors["text"]
            
            y_pos += 0.45
    
    def _add_stories_slide(self, slide, data: Dict, colors: Dict):
        """User stories slide with RICE scores"""
        stories = []
        if isinstance(data, dict) and 'stories' in data:
            stories = data['stories']
        elif isinstance(data, list):
            stories = data
        
        if not stories:
            self._add_text_box(slide, "", "No user stories available", Inches(0.5), Inches(1.5), colors)
            return
        
        # Table header
        rows = min(len(stories) + 1, 7)
        cols = 5
        
        table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(12), Inches(0.5 * rows)).table
        
        headers = ["ID", "User Story", "Priority", "Points", "RICE"]
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = colors["primary"]
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(11)
            paragraph.font.bold = True
            paragraph.font.color.rgb = colors["white"]
        
        # Set column widths
        table.columns[0].width = Inches(0.8)
        table.columns[1].width = Inches(6)
        table.columns[2].width = Inches(1.2)
        table.columns[3].width = Inches(1)
        table.columns[4].width = Inches(1)
        
        # Data rows
        for i, story in enumerate(stories[:6]):
            row = i + 1
            
            # ID
            table.cell(row, 0).text = str(story.get('id', i + 1))
            
            # Story description
            desc = story.get('description', '')
            if not desc and story.get('asA'):
                desc = f"As a {story.get('asA')}, I want {story.get('iWant', '...')} so that {story.get('soThen', '...')}"
            table.cell(row, 1).text = desc[:80]
            
            # Priority
            table.cell(row, 2).text = str(story.get('priority', '-'))
            
            # Story points
            table.cell(row, 3).text = str(story.get('storyPoints', '-'))
            
            # RICE score
            rice = story.get('riceScore', story.get('RICE', {}).get('score', '-'))
            table.cell(row, 4).text = str(rice)[:4]
            
            # Format cells
            for col in range(cols):
                cell = table.cell(row, col)
                cell.text_frame.paragraphs[0].font.size = Pt(9)
    
    def _add_roadmap_slide(self, slide, data: Dict, colors: Dict):
        """Roadmap slide with timeline"""
        phases = data.get('phases', [])
        if not phases and isinstance(data, list):
            phases = data
        
        if not phases:
            self._add_text_box(slide, "", "No roadmap data available", Inches(0.5), Inches(1.5), colors)
            return
        
        # Timeline visualization
        y_pos = 1.5
        
        for i, phase in enumerate(phases[:4]):
            # Phase card
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y_pos), Inches(12), Inches(1.2)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = colors["light"]
            card.line.color.rgb = colors["secondary"]
            
            # Phase number circle
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(0.8), Inches(y_pos + 0.2), Inches(0.6), Inches(0.6)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = colors["primary"]
            circle.line.fill.background()
            
            num_text = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos + 0.3), Inches(0.6), Inches(0.4))
            num_frame = num_text.text_frame
            num_frame.text = str(i + 1)
            num_frame.paragraphs[0].font.size = Pt(18)
            num_frame.paragraphs[0].font.bold = True
            num_frame.paragraphs[0].font.color.rgb = colors["white"]
            num_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Phase name
            name_box = slide.shapes.add_textbox(Inches(1.7), Inches(y_pos + 0.15), Inches(4), Inches(0.4))
            name_frame = name_box.text_frame
            name_frame.text = phase.get('name', f'Phase {i + 1}')
            name_frame.paragraphs[0].font.size = Pt(16)
            name_frame.paragraphs[0].font.bold = True
            name_frame.paragraphs[0].font.color.rgb = colors["primary"]
            
            # Duration
            if phase.get('duration'):
                dur_box = slide.shapes.add_textbox(Inches(1.7), Inches(y_pos + 0.55), Inches(3), Inches(0.3))
                dur_frame = dur_box.text_frame
                dur_frame.text = f"⏱ {phase['duration']}"
                dur_frame.paragraphs[0].font.size = Pt(11)
                dur_frame.paragraphs[0].font.color.rgb = colors["secondary"]
            
            # Milestones
            milestones = phase.get('milestones', [])
            if milestones:
                mile_box = slide.shapes.add_textbox(Inches(6), Inches(y_pos + 0.15), Inches(6), Inches(0.9))
                mile_frame = mile_box.text_frame
                mile_text = "Milestones:\n" + "\n".join([f"• {m[:40]}" for m in milestones[:3]])
                mile_frame.text = mile_text
                mile_frame.paragraphs[0].font.size = Pt(10)
                mile_frame.paragraphs[0].font.color.rgb = colors["text"]
            
            y_pos += 1.4
    
    def _add_okrs_slide(self, slide, data: Dict, colors: Dict, include_charts: bool):
        """OKRs slide with progress visualization"""
        okrs = [data.get('okr1'), data.get('okr2'), data.get('okr3')]
        okrs = [okr for okr in okrs if okr]
        
        if not okrs:
            self._add_text_box(slide, "", "No OKR data available", Inches(0.5), Inches(1.5), colors)
            return
        
        # North star
        if data.get('northStarDefinition'):
            self._add_info_card(slide, "🌟 North Star", data['northStarDefinition'], Inches(0.5), Inches(1.4), Inches(12), Inches(0.8), colors["accent"])
        
        y_pos = 2.4
        
        for i, okr in enumerate(okrs[:3]):
            # OKR card
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y_pos), Inches(12), Inches(1.3)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = colors["light"]
            card.line.color.rgb = colors["primary"]
            
            # Objective
            obj_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos + 0.1), Inches(11.5), Inches(0.4))
            obj_frame = obj_box.text_frame
            obj_frame.text = f"O{i+1}: {okr.get('objective', 'Objective')[:80]}"
            obj_frame.paragraphs[0].font.size = Pt(14)
            obj_frame.paragraphs[0].font.bold = True
            obj_frame.paragraphs[0].font.color.rgb = colors["primary"]
            
            # Key results
            krs = okr.get('keyResults', [])
            if krs:
                kr_text = "Key Results: " + " | ".join([f"KR{i+1}: {kr[:40]}" for i, kr in enumerate(krs[:3])])
                kr_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos + 0.6), Inches(11.5), Inches(0.6))
                kr_frame = kr_box.text_frame
                kr_frame.text = kr_text
                kr_frame.paragraphs[0].font.size = Pt(10)
                kr_frame.paragraphs[0].font.color.rgb = colors["text"]
                kr_frame.word_wrap = True
            
            y_pos += 1.5
    
    def _add_generic_slide(self, slide, data: Dict, colors: Dict):
        """Generic fallback slide"""
        self._add_text_box(slide, "Step Data:", str(data)[:500], Inches(0.5), Inches(1.5), colors)
    
    def _add_closing_slide(self, prs: Presentation, colors: Dict):
        """Thank you slide"""
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Background
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
            prs.slide_width, prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = colors["primary"]
        background.line.fill.background()
        
        # Thank you text
        thank_you = slide.shapes.add_textbox(Inches(0), Inches(2.5), prs.slide_width, Inches(1.5))
        thank_frame = thank_you.text_frame
        thank_frame.text = "Thank You"
        thank_frame.paragraphs[0].font.size = Pt(60)
        thank_frame.paragraphs[0].font.bold = True
        thank_frame.paragraphs[0].font.color.rgb = colors["white"]
        thank_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Subtitle
        sub = slide.shapes.add_textbox(Inches(0), Inches(4), prs.slide_width, Inches(0.8))
        sub_frame = sub.text_frame
        sub_frame.text = "Questions & Discussion"
        sub_frame.paragraphs[0].font.size = Pt(28)
        sub_frame.paragraphs[0].font.color.rgb = colors["light"]
        sub_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Product Pilot branding
        brand = slide.shapes.add_textbox(Inches(0), Inches(6.5), prs.slide_width, Inches(0.5))
        brand_frame = brand.text_frame
        brand_frame.text = "Generated by Product Pilot"
        brand_frame.paragraphs[0].font.size = Pt(14)
        brand_frame.paragraphs[0].font.color.rgb = colors["light"]
        brand_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Helper methods
    def _add_text_box(self, slide, label: str, text: str, x, y, colors: Dict, bold_label: bool = False):
        """Add a labeled text box"""
        if label:
            label_box = slide.shapes.add_textbox(x, y, Inches(12), Inches(0.3))
            label_frame = label_box.text_frame
            label_frame.text = label
            label_frame.paragraphs[0].font.size = Pt(11)
            label_frame.paragraphs[0].font.bold = True
            label_frame.paragraphs[0].font.color.rgb = colors["secondary"] if not bold_label else colors["primary"]
            y += Inches(0.35)
        
        text_box = slide.shapes.add_textbox(x, y, Inches(12), Inches(1))
        text_frame = text_box.text_frame
        text_frame.text = text
        text_frame.paragraphs[0].font.size = Pt(14)
        text_frame.paragraphs[0].font.color.rgb = colors["text"]
        text_frame.word_wrap = True
    
    def _add_wrapped_text(self, slide, text: str, x, y, width, font_size: Pt, color):
        """Add wrapped text"""
        box = slide.shapes.add_textbox(x, y, width, Inches(1))
        frame = box.text_frame
        frame.text = text
        frame.paragraphs[0].font.size = font_size
        frame.paragraphs[0].font.color.rgb = color
        frame.word_wrap = True
    
    def _add_bullet_list(self, slide, label: str, items: List[str], x, y, width, colors: Dict):
        """Add a bullet list"""
        if label:
            label_box = slide.shapes.add_textbox(x, y, width, Inches(0.3))
            label_frame = label_box.text_frame
            label_frame.text = label
            label_frame.paragraphs[0].font.size = Pt(11)
            label_frame.paragraphs[0].font.bold = True
            label_frame.paragraphs[0].font.color.rgb = colors["primary"]
            y += Inches(0.35)
        
        bullet_text = "\n".join([f"• {item[:60]}" for item in items[:5]])
        bullet_box = slide.shapes.add_textbox(x, y, width, Inches(2))
        bullet_frame = bullet_box.text_frame
        bullet_frame.text = bullet_text
        bullet_frame.paragraphs[0].font.size = Pt(11)
        bullet_frame.paragraphs[0].font.color.rgb = colors["text"]
    
    def _add_info_card(self, slide, title: str, content: str, x, y, width, height, accent_color):
        """Add an info card with accent color"""
        # Card background
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(249, 250, 251)
        card.line.color.rgb = accent_color
        
        # Title
        title_box = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.05), width - Inches(0.2), Inches(0.3))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(11)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = accent_color
        
        # Content
        content_box = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.35), width - Inches(0.2), height - Inches(0.4))
        content_frame = content_box.text_frame
        content_frame.text = content[:200]
        content_frame.paragraphs[0].font.size = Pt(10)
        content_frame.word_wrap = True
    
    def _get_date(self) -> str:
        """Get current date string"""
        from datetime import datetime
        return datetime.now().strftime("%B %d, %Y")
